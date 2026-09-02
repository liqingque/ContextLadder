#!/usr/bin/env python
"""Build the semi-final slide deck (PPTX, editable).

Covers the five items the submission specification requires: scientific problem,
method and technical route, results with settings and controls, reproduction
status, and open-source status. The specification also advises putting the most
important problem, capability and result in the first three pages, which is why
the finding leads rather than the architecture.

python-pptx is intentionally NOT installed into the frozen `tl` environment --
configs/env-spec.yaml and requirements.txt are hashed into
REPRODUCIBILITY_MANIFEST.json. Install it into an isolated directory and point
PYTHONPATH at it:

    pip install --target /tmp/pptxlib python-pptx
    PYTHONPATH=/tmp/pptxlib python scripts/build_slides.py --output slides.pptx
"""
import argparse
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Emu, Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "outputs" / "slides_assets"
NAVY, BLUE, GREY, RED = RGBColor(0x1A, 0x365, 0x60) if False else RGBColor(0x1A, 0x36, 0x5D), \
    RGBColor(0x2B, 0x6C, 0xB0), RGBColor(0x4A, 0x55, 0x68), RGBColor(0xC5, 0x30, 0x30)
CJK = "Noto Sans CJK SC"


def _tf(shape, size, color=GREY, bold=False):
    tf = shape.text_frame
    tf.word_wrap = True
    for p in tf.paragraphs:
        for r in p.runs:
            r.font.size, r.font.color.rgb, r.font.bold, r.font.name = Pt(size), color, bold, CJK
    return tf


def title_slide(prs, title, subtitle, foot):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    box = s.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(11.7), Inches(1.5))
    box.text_frame.text = title
    _tf(box, 40, NAVY, True)
    b2 = s.shapes.add_textbox(Inches(0.8), Inches(3.3), Inches(11.7), Inches(1.4))
    b2.text_frame.text = subtitle
    _tf(b2, 20, BLUE)
    b3 = s.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(11.7), Inches(0.8))
    b3.text_frame.text = foot
    _tf(b3, 13, GREY)
    return s


def content_slide(prs, title, bullets, image=None, img_top=2.7, img_h=3.4):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.9))
    t.text_frame.text = title
    _tf(t, 27, NAVY, True)
    if bullets:
        body = s.shapes.add_textbox(Inches(0.75), Inches(1.45), Inches(12), Inches(1.2 if image else 5))
        tf = body.text_frame
        tf.word_wrap = True
        for i, (txt, lvl, hi) in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = ("• " if lvl == 0 else "    – ") + txt
            p.space_after = Pt(6)
            for r in p.runs:
                r.font.size = Pt(16 if lvl == 0 else 14)
                r.font.color.rgb = RED if hi else GREY
                r.font.bold = bool(hi)
                r.font.name = CJK
    if image:
        pic = ASSETS / image
        if pic.exists():
            s.shapes.add_picture(str(pic), Inches(1.4), Inches(img_top), height=Inches(img_h))
    return s


def table_slide(prs, title, header, rows, note=None, col_w=None):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    t = s.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(12), Inches(0.9))
    t.text_frame.text = title
    _tf(t, 27, NAVY, True)
    n_r, n_c = len(rows) + 1, len(header)
    tbl = s.shapes.add_table(n_r, n_c, Inches(0.8), Inches(1.5),
                             Inches(11.7), Inches(0.42 * n_r)).table
    if col_w:
        total = sum(col_w)
        for i, w in enumerate(col_w):
            tbl.columns[i].width = Emu(int(Inches(11.7) * w / total))
    for j, h in enumerate(header):
        c = tbl.cell(0, j); c.text = h
        for p in c.text_frame.paragraphs:
            for r in p.runs:
                r.font.size, r.font.bold, r.font.name = Pt(13), True, CJK
                r.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    for i, row in enumerate(rows, start=1):
        for j, v in enumerate(row):
            c = tbl.cell(i, j); c.text = str(v)
            for p in c.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size, r.font.name = Pt(12), CJK
                    r.font.color.rgb = RED if str(v).startswith("★") else GREY
    if note:
        b = s.shapes.add_textbox(Inches(0.8), Inches(1.6 + 0.42 * n_r), Inches(11.7), Inches(1.0))
        b.text_frame.text = note
        _tf(b, 13, GREY)
    return s


def build(out):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    title_slide(prs, "上下文阶梯 ContextLadder",
                "面向未见扰动与未见菌株的虚拟酵母蛋白组预测\n\n"
                "核心结论：占官方权重 20% 的模块上，零样本化合物泛化目前尚未被真正测量到",
                "GOAI 赛道三 · 虚拟细胞方向 ｜ 团队 TorchDragon（李晓蒙、徐逸飞）｜ 2026-09-01")

    content_slide(prs, "1. 科学问题：这个任务的分数，具体在测什么？", [
        ("官方 20% 的「上下文均值残差」是最直接对应零样本扰动泛化的一项", 0, False),
        ("我们先测了它：同一上下文内、携带不同未见化合物的行对，预测向量差多少", 0, False),
        ("S1：379 个上下文、991 个跨化合物行对，991 对逐位相同，max|Δ| = 0", 1, True),
        ("S3：96 个上下文、250 对，250 对逐位相同，max|Δ| = 0", 1, True),
        ("五个候选（含两类官方基线家族构造）结果完全一致；更早八个架构同样如此", 0, False),
        ("→ 该模块的得分完全来自所有化合物共享的分量，没有化合物特异信息参与", 0, True),
    ])

    content_slide(prs, "2. 系统能力：官方六模块实测", [
        ("在化合物信息确实可用的场景上，本方案是平凡基线的数倍", 0, False),
        ("时间插值 FC 0.6717 vs 平凡基线 0.2291 —— 2.9 倍，且是全部场景最高", 1, True),
        ("DEP 方向准确率 0.8919 vs 0.7285；这 15% 权重此前从未被任何 gate 评过", 1, False),
    ], image="modules.png", img_top=2.85, img_h=3.3)

    table_slide(prs, "3. 关键结果：与官方基线家族的对照",
                ["候选", "W", "ΔW", "裁决"],
                [["本方案 ContextLadder", "0.5218", "—", "基线"],
                 ["分场景低秩统计", "0.4471", "−0.0747", "DOMINATED"],
                 ["分场景梯度提升", "0.4424", "−0.0794", "DOMINATED"],
                 ["全局梯度提升 PCA-128", "0.4186", "−0.1032", "DOMINATED"],
                 ["全局梯度提升 PCA-64", "0.4181", "−0.1037", "DOMINATED"]],
                note="采纳门槛 ΔW ≥ +0.010（= 已冻结的 S1 效应量地板 0.053 × 该模块 0.20 权重）。"
                     "实测为反方向的 7.5×~10.4×，六个加权模块无一例外全部更差，S1 的四个差值经配对聚类 "
                     "bootstrap 全部 CI 不含 0。\n限定：三种构造一天内实现，GBDT 迭代上限 25/20，两种夹了 "
                     "PCA 瓶颈——可主张的是「这三种具体构造被劣于」，而非「梯度提升走不通」。",
                col_w=[4, 2, 2, 3])

    content_slide(prs, "4. 选择目标 ≠ 评分目标：两个会影响所有队伍的口径问题", [
        ("对照池从哪里取，赛题未限定，而它会静默改变被评估的样本集合", 0, False),
        ("官方口径 2,806 匹配行；改用仅训练池对照会丢掉 1,602 行 = 57.1%", 1, True),
        ("丢掉的恰好是 S2 与 S3 两个最难、得分最低的场景", 1, False),
        ("我们自己 08-13 至 08-15 的多轮 gate 就踩在这个坑里", 1, False),
        ("这个划分能验证什么、不能验证什么", 0, False),
        ("S1 仅 6 个聚类 → 效应量地板 0.053；小于它的改动在本划分上不可判定", 1, True),
        ("S2 仅 1 个聚类（BAI），而 test 主力未见菌株 CRD 在 train_val 中零覆盖", 1, True),
        ("→ 选择实体与评测实体不同，整条菌株侧建模分支据此关闭", 1, False),
    ])

    content_slide(prs, "5. 容量-内容分离：一条给社区的定量路标", [
        ("嵌套无泄漏 37 折 LOCO，判据在数据进入前冻结", 0, False),
        ("容量够：全张成 oracle ρ = 0.2671，单侧 95% 上界 0.2941 / 0.2706 > 0.20 门槛", 1, False),
        ("解码器不是瓶颈：稠密 − 稀疏在四个通道上全部为负", 1, False),
        ("缺的是内容：所有可构造通道都打不过负对照，maxT 校正后同样不显著", 1, True),
        ("稀疏先验族数学上关闭：符号最优上限 0.0491 < 效应量地板 0.053", 1, True),
    ], image="capacity_content.png", img_top=3.35, img_h=2.9)

    content_slide(prs, "6. 方法：ContextLadder 是上述测量的推论，不是零件拼装", [
        ("约束：对未见化合物没有可用的化合物特异信息 → 唯一有证据的行为是退回更粗一级", 0, True),
        ("测量分支的分层收缩 —— 先有测量，后有设计", 0, False),
        ("measurement-only Ridge 0.982272 vs biological-only 0.962347 → batch shortcut 风险", 1, False),
        ("把板号建模为来源/仪器之上的收缩残差，未见板结构性回退", 1, False),
        ("化合物掩码 —— 把「不知道」从随机向量变成学出来的状态", 0, False),
        ("事后均值替换方向对但不稳健（一个种子 −0.0124，另两个 ≈0；菌株侧全退化）", 1, False),
        ("训练式掩码三种子全部改善，集成 RMSE −0.003405，两个 OOD 子集同时改善", 1, False),
        ("双字段掩码是负结果（+0.015）→ 设计有方向性、可被证伪，不是随手加的正则", 1, True),
        ("创新性主张的边界：不主张 FiLM / 掩码 / 集成本身新颖，主张把「如何优雅地不知道」当作一等设计目标", 0, False),
    ])

    table_slide(prs, "7. 赛题四场景 ↔ 机制 ↔ 实测（官方模块与权重）",
                ["赛题场景", "官方模块 / 权重", "本方案机制", "验证集实测"],
                [["未见化合物", "上下文均值残差 S1 ｜ 20%",
                  "25% 概率掩码化合物 token 训练，未见化合物落到学出的 <UNK> 回退态",
                  "S1 0.2034；化合物特异信号为零是全体候选共性"],
                 ["未见菌株", "药物均值残差 S2 ｜ 20%",
                  "UNK 槽位 + 共享生物/测量上下文；菌株侧改进分支因证据不足整族关闭",
                  "S2 0.4077（单聚类，仅描述性）"],
                 ["双重未知", "匹配对照 FC 全体 ｜ 25%",
                  "两个回退同时生效，无未见实体专属参数被激活",
                  "含于 FC 0.4914；S3 行对 250 对逐位相同"],
                 ["时间插值", "时间插值 FC ｜ 10%",
                  "时间为连续协变量（z-score + RBF 基），任意时刻自然插值",
                  "0.6717 = 平凡基线 2.9×，全场景最高"]],
                note="四个场景是同一回退原则的四次实例化：宁可退到更粗的一级，也不要在没有证据的地方假装知道。\n"
                     "数字与报告第 5.5 / 6 节同源（冻结验证集，对照池 train | val）。",
                col_w=[2, 3, 5, 4])

    table_slide(prs, "8. 证伪阶梯：十五族分支，判据先于数据",
                ["#", "分支族", "代表实测", "裁决"],
                [["2", "化学结构特征", "6 个未见化合物 ΔRMSE +0.097757", "拒绝"],
                 ["3", "基因组特征", "genome kernel LOSO +0.025328", "关联 ≠ 预测"],
                 ["4", "外部蛋白组迁移", "最差菌株 CGD +0.087488", "系数不可识别"],
                 ["5", "逐蛋白后处理校正", "EB 校准 +0.004050", "家族级关闭"],
                 ["7", "稀疏机制先验", "符号最优上限 0.0491 < 0.053", "数学上关闭"],
                 ["9", "优化侧杠杆", "最优 ΔW +0.0075；地板重标定后见第 14 行", "NOT ADOPTED"],
                 ["12", "容量 / EMA / 跨架构混合", "十候选最大 ΔW +0.0050", "NOT ADOPTED"],
                 ["13", "低速率菌株掩码", "0.05 → +0.0037，0.10 → +0.0119", "假说被证伪"],
                 ["14", "跨家族组合", "★通过五条采纳条件；换族翻符号，十配对 3/10 改善", "采纳后被推翻"]],
                note="★ 三条自证记录：(a) 第 8 族若把 δ 定为 0.10 本轮即满足 REPRESENTATION GAP，"
                     "但 δ 在看到任何 ρ 之前已冻结为 0.05，如实判 INDETERMINATE；"
                     "(b) R1 凸混合是八个候选中唯一改善 20% 权重模块的（+0.02191，CI 不含 0），"
                     "仍按 0.053 地板判 FAIL——后证明是对的：纯种子噪声可给出几乎重合的 CI；"
                     "(c) 第 14 族的候选通过了全部五条冻结条件，被一次跨种子族复核推翻；"
                     "十种子配对比较确证为抽样假象（3/10 改善，配对差均值 +0.001418 即更差，符号检验 p = 0.344）。"
                     "推论：检测地板必须跨越所有未被测量的自由度。",
                col_w=[1, 4, 5, 3])

    table_slide(prs, "9. 复现情况：干净目录实测，逐字节可复现",
                ["核验项", "结果"],
                [["训练（三个种子，单张 RTX 3090）", "98.0 s，峰值显存 923 MiB"],
                 ["推理", "17.5 s"],
                 ["全新 venv 装 requirements", "exit 0，新环境内 pytest 12 passed"],
                 ["三个 checkpoint SHA-256", "★与随包 manifest 完全一致"],
                 ["prediction.csv SHA-256", "★与随包文件同一值"],
                 ["提交格式校验", "PASS，九项硬检查全过"],
                 ["免数据契约测试", "12 passed"]],
                note="数字出处：reports/P0_REPRODUCTION_HARDENING_20260901.md（2026-09-01 干净目录实测）。\n"
                     "测试蛋白组从未放入干净目录，四条命令仍全部成功——predict.py 中不存在任何蛋白组读取路径，"
                     "测试对此有静态断言。\n已知限制：官方未下发 4,422 列 contract 文件，按规则实算并落盘列表与 "
                     "SHA-256；GPU 归约不保证跨硬件位确定。",
                col_w=[5, 7])

    content_slide(prs, "10. 开源情况", [
        ("仓库与许可：代码、权重与全部清单随提交包分发（逐文件 SHA-256 冻结，不依赖外部仓库可达）；赛后开源，许可 CC BY-NC 4.0", 0, False),
        ("最终模型不使用任何外部数据，只消费官方样本元数据字段", 0, False),
        ("外部资源只出现在证伪实验中：source_manifest.json（来源/版本/许可/校验和）+ entity_mapping.csv（65 行，含未匹配/歧义/代理）", 0, False),
        ("四件可被社区直接接续的工具", 0, True),
        ("官方口径六模块评估器 —— 任何队伍都可以用它自查那个 57% 的对照池问题", 1, False),
        ("化合物不变性检查 —— 替代零功效的置换检验，新候选应先过这一关再谈 S1 分数", 1, False),
        ("嵌套无泄漏 LOCO harness —— 37 折，逐折泄漏断言，独立实现逐位复现", 1, False),
        ("交叉拟合 OOF 基础设施与泄漏审计 —— 在此发现并修复了一个历史 both-fold 泄漏 bug", 1, False),
    ])

    Path(out).parent.mkdir(parents=True, exist_ok=True)
    prs.save(out)
    print(f"wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")


if __name__ == "__main__":
    ap = argparse.ArgumentParser()
    ap.add_argument("--output", default=str(ROOT / "outputs" / "slides.pptx"))
    build(ap.parse_args().output)
